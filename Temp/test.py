from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title_text, content_text, bullet_points=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    tf = content.text_frame
    tf.text = content_text

    if bullet_points:
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Problem 9.20 Solution"
    subtitle.text = "Transfer Function and Noise Analysis of First-Order Active Filters"

    # Slide 2: Part (a) - Circuit I Transfer Function
    title_s2 = "Part (a): Transfer Function of Circuit I"
    content_s2 = "Circuit I is a standard inverting integrator with a parallel feedback resistor."
    bullets_s2 = [
        "Feedback resistor R_F = 7 kΩ, Input resistor R_in = 7 kΩ",
        "Feedback capacitor C_F = 80 nF",
        "The transfer function is given by:",
        "Vo / Vi = -1 / (1 + jω * R_F * C_F)",
        "Substituting the values:",
        "Vo / Vi = -1 / (1 + jω * (7k) * (80n))  --- (1)"
    ]
    add_slide(prs, title_s2, content_s2, bullets_s2)

    # Slide 3: Part (a) - Circuit II Transfer Function
    title_s3 = "Part (a): Transfer Function of Circuit II"
    content_s3 = "Circuit II is a two-stage circuit. Op-amp 2 is a non-inverting amplifier with a gain of 2."
    bullets_s3 = [
        "Let the output of op-amp 1 be Vx. Then Vo = 2 * Vx, or Vx = Vo / 2.",
        "Apply KCL at the inverting terminal of op-amp 1 (virtual ground):",
        "Current through input resistor + Current through feedback resistor + Current through capacitor = 0",
        "(Vi / 14k) + (Vo / 14k) + (Vx / (1 / (jω * C_F))) = 0",
        "Substitute Vx = Vo / 2 and C_F = 80nF:",
        "(Vi / 14k) + (Vo / 14k) + (Vo / 2) * (jω * 80n) = 0"
    ]
    add_slide(prs, title_s3, content_s3, bullets_s3)

    # Slide 4: Part (a) - Conclusion
    title_s4 = "Part (a): Conclusion"
    content_s4 = "Simplifying the KCL equation for Circuit II."
    bullets_s4 = [
        "Multiply the equation by 14k:",
        "Vi + Vo + Vo * (jω * 80n * 7k) = 0",
        "Vi + Vo * (1 + jω * (7k) * (80n)) = 0",
        "Rearranging for the transfer function:",
        "Vo / Vi = -1 / (1 + jω * (7k) * (80n))  --- (2)",
        "Comparing (1) and (2), both circuits have identical transfer functions."
    ]
    add_slide(prs, title_s4, content_s4, bullets_s4)

    # Slide 5: Part (b) - Output Noise for Circuit I
    title_s5 = "Part (b): Output Noise for Circuit I"
    content_s5 = "Calculating total output noise using dominant noise sources (resistors)."
    bullets_s5 = [
        "Output noise spectral density at DC (from solution):",
        "V_no^2(0) = [I_n1^2 + I_nF^2] * R_F^2 = (15.2 nV/√Hz)^2",
        "Total RMS output noise for a single-pole low-pass filter:",
        "V_no(rms)^2 = V_no^2(0) * [1 / (4 * R_F * C_F)]",
        "Calculating the RMS value:",
        "V_no(rms) = V_no(0) / sqrt(4 * R_F * C_F)",
        "From solution: V_no(rms) = 0.32 μV"
    ]
    add_slide(prs, title_s5, content_s5, bullets_s5)

    # Slide 6: Part (c) - Output Noise for Circuit II (Main)
    title_s6 = "Part (c): Output Noise for Circuit II (Resistors)"
    content_s6 = "Calculating noise from the 14 kΩ resistors (R1 and R2)."
    bullets_s6 = [
        "Output noise spectral density at DC is higher due to larger resistors:",
        "V_no1^2(0) = [I_n1^2 + I_nF^2] * R_F^2 = (21.5 nV/√Hz)^2",
        "The effective noise bandwidth is determined by R_eff = R_F / 2 = 7k.",
        "Total RMS output noise:",
        "V_no1(rms)^2 = V_no1^2(0) * [1 / (4 * (R_F/2) * C_F)]",
        "From solution: V_no1(rms) = 0.45 μV"
    ]
    add_slide(prs, title_s6, content_s6, bullets_s6)

    # Slide 7: Part (c) - Noise from Op-amp 2
    title_s7 = "Part (c): Noise Current at Op-amp 2 Input"
    content_s7 = "Analyzing the effect of a noise current source (I_i2) at the negative terminal of op-amp 2."
    bullets_s7 = [
        "The transfer function from I_i2 to Vo is derived as:",
        "Vo / I_i2 = -R_2 / (1 + 1 / (jω * C_F * R_F / 2))",
        "Vo / I_i2 = (-jω * R_2 * R_F * C_F / 2) / (1 + jω * C_F * R_F / 2)",
        "This is a high-pass filter transfer function.",
        "Conclusion: For ideal white noise, the total output noise due to current noise at this terminal is infinite (∞), as high frequencies are not attenuated."
    ]
    add_slide(prs, title_s7, content_s7, bullets_s7)

    prs.save('Problem_9.20_Solution.pptx')

if __name__ == "__main__":
    create_presentation()